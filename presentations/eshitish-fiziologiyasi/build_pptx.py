# -*- coding: utf-8 -*-
"""
Eshitish fiziologiyasi — .pptx ni yigʻadi (python-pptx).
content.py dan matn, assets/ dan diagrammalar olinadi.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import content as C

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ---- Tema ----
INK    = RGBColor(0x0F, 0x2B, 0x36)
TEAL   = RGBColor(0x1F, 0x7A, 0x8C)
TEAL_D = RGBColor(0x13, 0x52, 0x5E)
TEAL_L = RGBColor(0xCF, 0xE5, 0xEA)
CREAM  = RGBColor(0xF6, 0xEF, 0xE2)
SAND   = RGBColor(0xEA, 0xD8, 0xB8)
ACCENT = RGBColor(0xE0, 0x7A, 0x5F)
AMBER  = RGBColor(0xE9, 0xB4, 0x4C)
ROSE   = RGBColor(0xD3, 0x6B, 0x8B)
GREEN  = RGBColor(0x5B, 0x8C, 0x5A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT   = RGBColor(0x16, 0x2A, 0x31)
GREY   = RGBColor(0x5A, 0x6A, 0x70)
PANEL  = RGBColor(0xFB, 0xF8, 0xF1)

FONT = "DejaVu Sans"  # oʻ/gʻ kabi belgilar uchun ishonchli

EMU = 914400
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf


def setpara(p, size, color, bold=False, align=PP_ALIGN.LEFT, space_after=6,
            italic=False, font=FONT, line=None):
    p.alignment = align
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic; r.font.name = font


def add_run(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return r


def title_band(s, title, idx):
    rect(s, 0, 0, SW, Inches(1.15), TEAL)
    rect(s, 0, Inches(1.15), SW, Inches(0.07), AMBER)
    rect(s, Inches(0.0), 0, Inches(0.18), Inches(1.15), AMBER)
    tb, tf = textbox(s, Inches(0.55), Inches(0.12), Inches(11.5), Inches(0.95),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; add_run(p, title, 30, WHITE, bold=True)
    setpara(p, 30, WHITE, bold=True)
    page_number(s, idx)


def page_number(s, idx):
    tb, tf = textbox(s, Inches(12.3), Inches(7.02), Inches(0.9), Inches(0.4))
    p = tf.paragraphs[0]; add_run(p, str(idx), 12, GREY)
    setpara(p, 12, GREY, align=PP_ALIGN.RIGHT, space_after=0)


def footer_note(s, text):
    rect(s, Inches(0.55), Inches(6.55), Inches(12.23), Inches(0.62), CREAM,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, Inches(0.75), Inches(6.58), Inches(11.8), Inches(0.55),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; add_run(p, text, 14, TEAL_D, italic=True)
    setpara(p, 14, TEAL_D, italic=True, space_after=0)


def pic(s, name, x, y, w=None, h=None):
    path = os.path.join(ASSETS, name + ".png")
    if not os.path.exists(path):
        return None
    return s.shapes.add_picture(path, x, y, width=w, height=h)


def bullet_para(tf, term, desc, first=False, size=17, gap=10):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    add_run(p, "▸  ", size, ACCENT, bold=True)
    if term:
        add_run(p, term + (" — " if desc else ""), size, INK, bold=True)
    if desc:
        add_run(p, desc, size, TEXT)
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(gap)
    p.line_spacing = 1.08
    return p


# ===================================================================
# Slayd qurувchilar
# ===================================================================
def build_title(s, d, idx):
    bg(s, INK)
    # hero rasm (oʻng tomon to'liq)
    if os.path.exists(os.path.join(ASSETS, "hero.png")):
        pic(s, "hero", Inches(7.0), 0, w=Inches(6.333), h=SH)
    rect(s, 0, 0, Inches(7.4), SH, INK)
    rect(s, Inches(0.0), Inches(2.2), Inches(0.25), Inches(2.6), AMBER)
    tb, tf = textbox(s, Inches(0.7), Inches(2.1), Inches(6.6), Inches(3.4))
    p = tf.paragraphs[0]; add_run(p, d["title"], 50, WHITE, bold=True)
    setpara(p, 50, WHITE, bold=True, space_after=14, line=1.04)
    p2 = tf.add_paragraph(); add_run(p2, d["subtitle"], 21, TEAL_L)
    setpara(p2, 21, TEAL_L, space_after=10, line=1.1)
    tb2, tf2 = textbox(s, Inches(0.7), Inches(6.2), Inches(6.0), Inches(1.0))
    for i, m in enumerate(d.get("meta", [])):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        add_run(p, m, 15, RGBColor(0xAD, 0xC4, 0xCB))
        setpara(p, 15, RGBColor(0xAD, 0xC4, 0xCB), space_after=2)


def build_agenda(s, d, idx):
    bg(s, PANEL)
    title_band(s, d["title"], idx)
    items = d["items"]
    n = len(items)
    top, bottom = 1.55, 6.95
    row = (bottom - top) / n
    dia = 0.52
    for i, item in enumerate(items):
        cy = top + row * i
        # raqamli doiracha
        c = rect(s, Inches(0.8), Inches(cy + (row - dia) / 2), Inches(dia), Inches(dia),
                 TEAL, shape=MSO_SHAPE.OVAL)
        ctf = c.text_frame; ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; add_run(cp, str(i + 1), 17, WHITE, bold=True)
        setpara(cp, 17, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        # matn (doiracha bilan bir qatorda, vertikal markazda)
        tb, tf = textbox(s, Inches(1.6), Inches(cy), Inches(11.2), Inches(row),
                         anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; add_run(p, item, 18, TEXT)
        setpara(p, 18, TEXT, space_after=0, line=1.0)


def build_bullets(s, d, idx):
    bg(s, PANEL)
    title_band(s, d["title"], idx)
    tb, tf = textbox(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(5.4))
    for i, (term, desc) in enumerate(d["bullets"]):
        bullet_para(tf, term, desc, first=(i == 0), size=19, gap=14)


def build_image_bullets(s, d, idx):
    bg(s, PANEL)
    title_band(s, d["title"], idx)
    # chap: bullets (rasm panelidan aniq ajratilgan)
    tb, tf = textbox(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(5.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    for i, b in enumerate(d["bullets"]):
        if isinstance(b, (list, tuple)):
            term, desc = b
        else:
            term, desc = "", b
        bullet_para(tf, term, desc, first=(i == 0), size=15, gap=10)
    # o'ng: rasm panel
    rect(s, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.0), WHITE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL_L)
    p = pic(s, d["image"], Inches(7.05), Inches(1.65))
    if p:
        # rasmni panelga moslab o'lchamlash
        maxw, maxh = Inches(5.6), Inches(4.7)
        ratio = p.width / p.height
        if p.width > maxw:
            p.width = maxw; p.height = Emu(int(maxw / ratio))
        if p.height > maxh:
            p.height = maxh; p.width = Emu(int(maxh * ratio))
        p.left = Inches(6.9) + Emu(int((Inches(5.9) - p.width) / 2))
        p.top = Inches(1.5) + Emu(int((Inches(5.0) - p.height) / 2))
    if d.get("footer"):
        footer_note(s, d["footer"])


def build_steps(s, d, idx):
    bg(s, PANEL)
    title_band(s, d["title"], idx)
    steps = d["steps"]
    n = len(steps)
    cols = [TEAL_D, TEAL, RGBColor(0x2E, 0x86, 0x95), AMBER, ACCENT, ROSE, GREEN]
    y0 = Inches(1.7); bh = Inches(0.62); gap = Inches(0.07)
    x = Inches(0.8); w = Inches(11.7)
    for i, st in enumerate(steps):
        y = y0 + (bh + gap) * i
        c = cols[i % len(cols)]
        tc = INK if c == AMBER else WHITE
        r = rect(s, x, y, Inches(0.62), bh, c, shape=MSO_SHAPE.OVAL)
        rp = r.text_frame.paragraphs[0]; add_run(rp, str(i + 1), 16, tc, bold=True)
        setpara(rp, 16, tc, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        bar = rect(s, x + Inches(0.78), y, w - Inches(0.78), bh, WHITE,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL_L)
        btf = bar.text_frame; btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; add_run(bp, st, 16, TEXT, bold=(i in (5,)))
        setpara(bp, 16, TEXT, align=PP_ALIGN.LEFT, space_after=0)
        bp.runs[0].font.bold = False
    if d.get("footer"):
        footer_note(s, d["footer"])


def build_two_col(s, d, idx):
    bg(s, PANEL)
    title_band(s, d["title"], idx)

    def column(x, head, items, col):
        rect(s, x, Inches(1.6), Inches(5.9), Inches(0.8), col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb, tf = textbox(s, x, Inches(1.62), Inches(5.9), Inches(0.76),
                         anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; add_run(p, head, 18, WHITE, bold=True)
        setpara(p, 18, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        rect(s, x, Inches(2.5), Inches(5.9), Inches(4.4), WHITE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL_L)
        tb2, tf2 = textbox(s, x + Inches(0.25), Inches(2.7), Inches(5.4), Inches(4.0))
        for i, it in enumerate(items):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            add_run(p, "•  ", 17, col, bold=True)
            add_run(p, it, 17, TEXT)
            p.space_after = Pt(12); p.line_spacing = 1.1

    column(Inches(0.55), d["left_title"], d["left"], TEAL_D)
    column(Inches(6.9), d["right_title"], d["right"], ACCENT)


def build_summary(s, d, idx):
    bg(s, PANEL)
    title_band(s, d["title"], idx)
    tb, tf = textbox(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(3.5))
    for i, pt in enumerate(d["points"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        add_run(p, "✓  ", 18, GREEN, bold=True)
        add_run(p, pt, 18, TEXT)
        p.space_after = Pt(13); p.line_spacing = 1.12
    # manbalar paneli
    rect(s, Inches(0.7), Inches(5.15), Inches(11.93), Inches(1.95), CREAM,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb2, tf2 = textbox(s, Inches(1.0), Inches(5.25), Inches(11.4), Inches(1.8))
    p = tf2.paragraphs[0]; add_run(p, "Manbalar", 16, TEAL_D, bold=True)
    setpara(p, 16, TEAL_D, bold=True, space_after=4)
    for src in d["sources"]:
        pp = tf2.add_paragraph(); add_run(pp, "• " + src, 13, GREY)
        pp.space_after = Pt(2); pp.line_spacing = 1.05


BUILDERS = {
    "title": build_title,
    "agenda": build_agenda,
    "bullets": build_bullets,
    "image_bullets": build_image_bullets,
    "steps": build_steps,
    "two_col": build_two_col,
    "summary": build_summary,
}


def main():
    for i, d in enumerate(C.SLIDES, start=1):
        s = slide()
        BUILDERS[d["kind"]](s, d, i)
    out = os.path.join(HERE, "Eshitish_fiziologiyasi.pptx")
    prs.save(out)
    print("saved", out, "| slaydlar:", len(C.SLIDES))


if __name__ == "__main__":
    main()
