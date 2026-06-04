# -*- coding: utf-8 -*-
"""
Tekshirish uchun: haqiqiy .pptx fayldagi shakllarni oʻqib, har bir slaydni
PNG koʻrinishida chizadi (matnlar, ranglar, rasm joylashuvi). LibreOffice
ishlamagani uchun shu yordamchi orqali kompozitsiyani koʻrish mumkin.
Bu rasmlar repozitoriyga qoʻshilmaydi — faqat /tmp ga chiqadi.
"""
import io
import os
import sys
import textwrap
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, "Eshitish_fiziologiyasi.pptx")
OUT = "/tmp/preview"
os.makedirs(OUT, exist_ok=True)

EMU_PER_IN = 914400
prs = Presentation(PPTX)
SW = prs.slide_width / EMU_PER_IN
SH = prs.slide_height / EMU_PER_IN

ALIGN = {PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center",
         PP_ALIGN.RIGHT: "right", None: "left"}


def rgb_of(color, default=None):
    try:
        c = color.rgb
        return "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        return default


def shape_fill(sp):
    try:
        if sp.fill.type is not None:
            return rgb_of(sp.fill.fore_color, None)
    except Exception:
        pass
    return None


def draw_text(ax, sp, x, y, w, h):
    tf = sp.text_frame
    # vertikal tayanch
    va = getattr(tf, "vertical_anchor", None)
    runs_all = []
    for p in tf.paragraphs:
        txt = "".join(r.text for r in p.runs) or p.text
        if not txt:
            runs_all.append(("", 14, "#000000", False, "left"))
            continue
        size = 14; color = "#16242b"; bold = False
        for r in p.runs:
            if r.font.size:
                size = max(size, r.font.size.pt)
            rc = rgb_of(r.font.color, None)
            if rc:
                color = rc
            if r.font.bold:
                bold = True
        al = ALIGN.get(p.alignment, "left")
        runs_all.append((txt, size, color, bold, al))

    # qatorlarni oʻrash (overflow ni koʻrish uchun)
    lines = []
    for txt, size, color, bold, al in runs_all:
        maxchars = max(4, int(w * 132 / size))
        wrapped = textwrap.wrap(txt, maxchars) or [""]
        for j, ln in enumerate(wrapped):
            lines.append((ln, size, color, bold, al))

    total_h = sum(s * 1.45 / 72 for _, s, _, _, _ in lines)
    # y oʻqi teskari (yuqori = kichik y), pastga qarab cy ni OSHIRAMIZ
    if va == MSO_ANCHOR.MIDDLE:
        cy = y + h / 2 - total_h / 2
    elif va == MSO_ANCHOR.BOTTOM:
        cy = y + h - total_h
    else:
        cy = y + 0.06
    overflow = total_h > h + 0.06
    for ln, size, color, bold, al in lines:
        if al == "center":
            tx = x + w / 2
        elif al == "right":
            tx = x + w - 0.05
        else:
            tx = x + 0.06
        ax.text(tx, cy, ln, fontsize=size * 0.96, color=color,
                ha=al, va="top", fontweight="bold" if bold else "normal",
                family="DejaVu Sans")
        cy += size * 1.45 / 72
    return overflow


def render_slide(slide, idx):
    fig, ax = plt.subplots(figsize=(SW, SH), dpi=110)
    ax.set_xlim(0, SW); ax.set_ylim(0, SH); ax.axis("off")
    ax.invert_yaxis()
    bgc = None
    try:
        bgc = rgb_of(slide.background.fill.fore_color, None)
    except Exception:
        pass
    ax.add_patch(Rectangle((0, 0), SW, SH, fc=bgc or "#ffffff", zorder=0))
    overflows = []
    for sp in slide.shapes:
        x = sp.left / EMU_PER_IN if sp.left is not None else 0
        y = sp.top / EMU_PER_IN if sp.top is not None else 0
        w = sp.width / EMU_PER_IN if sp.width is not None else 1
        h = sp.height / EMU_PER_IN if sp.height is not None else 1
        if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = Image.open(io.BytesIO(sp.image.blob)).convert("RGBA")
                ax.imshow(img, extent=(x, x + w, y + h, y), zorder=3,
                          aspect="auto")
            except Exception as e:
                ax.add_patch(Rectangle((x, y), w, h, fc="#dddddd", zorder=3))
        else:
            fc = shape_fill(sp)
            if fc:
                ax.add_patch(FancyBboxPatch((x, y), w, h,
                             boxstyle="round,pad=0,rounding_size=0.05",
                             fc=fc, ec="none", zorder=2))
            if sp.has_text_frame and sp.text_frame.text.strip():
                of = draw_text(ax, sp, x, y, w, h)
                if of:
                    overflows.append(sp.text_frame.text[:40])
    ax.text(SW - 0.1, 0.18, f"#{idx}", ha="right", va="top",
            fontsize=8, color="#999999")
    path = os.path.join(OUT, f"slide_{idx:02d}.png")
    fig.savefig(path, dpi=110, bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    return overflows


def main():
    only = [int(a) for a in sys.argv[1:]] if len(sys.argv) > 1 else None
    flagged = {}
    for i, s in enumerate(prs.slides, 1):
        if only and i not in only:
            continue
        of = render_slide(s, i)
        if of:
            flagged[i] = of
    print("rendered ->", OUT)
    if flagged:
        print("POTENTIAL OVERFLOW:")
        for k, v in flagged.items():
            print(f"  slide {k}: {v}")
    else:
        print("no overflow detected")


if __name__ == "__main__":
    main()
